"""
Extract images and video URLs from documents during indexing.
Images are uploaded to MinIO and referenced by URL.
"""
import io
import os
import re
from typing import Any, Dict, List

import boto3
from botocore.config import Config

# ── Video URL patterns ────────────────────────────────────────────────────────

_VIDEO_PATTERNS = [
    r'https?://(?:www\.)?youtube\.com/watch\?v=[\w-]+',
    r'https?://(?:www\.)?youtube\.com/shorts/[\w-]+',
    r'https?://youtu\.be/[\w-]+',
    r'https?://(?:www\.)?vimeo\.com/\d+',
    r'https?://(?:www\.)?loom\.com/share/[\w-]+',
    r'https?://(?:www\.)?drive\.google\.com/file/d/[\w-]+',
]
_VIDEO_RE = re.compile('|'.join(_VIDEO_PATTERNS), re.IGNORECASE)

# ── S3/MinIO helpers ──────────────────────────────────────────────────────────

def _s3():
    return boto3.client(
        "s3",
        endpoint_url=os.getenv("AWS_ENDPOINT_URL", "http://localhost:9000"),
        aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID", "staffbot"),
        aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY", "staffbot123"),
        region_name=os.getenv("AWS_REGION", "us-east-1"),
        config=Config(signature_version="s3v4"),
    )


def _upload_image(img_bytes: bytes, ext: str, tenant_id: str, document_id: str, index: int) -> str:
    """Upload image to MinIO, return public URL."""
    bucket = os.getenv("AWS_S3_BUCKET", "staffbot-docs")
    key    = f"{tenant_id}/{document_id}/images/img_{index}.{ext}"
    _s3().put_object(Bucket=bucket, Key=key, Body=img_bytes, ContentType=f"image/{ext}")
    # MINIO_PUBLIC_URL points to the nginx proxy which already maps to the bucket root,
    # so the URL is {public_base}/{key} (no bucket name).
    # Without it, fall back to direct MinIO: {endpoint}/{bucket}/{key}.
    public_base = os.getenv("MINIO_PUBLIC_URL", "").rstrip("/")
    if public_base:
        return f"{public_base}/{key}"
    endpoint = os.getenv("AWS_ENDPOINT_URL", "http://localhost:9000")
    return f"{endpoint}/{bucket}/{key}"


# ── Per-format extractors ─────────────────────────────────────────────────────

MAX_IMAGES_PER_DOCUMENT = 50   # cap to avoid unbounded MinIO uploads blocking the event loop

def _extract_pdf(content: bytes, tenant_id: str, document_id: str, upload_images: bool = True) -> Dict[str, Any]:
    import fitz  # pymupdf

    doc        = fitz.open(stream=content, filetype="pdf")
    pages_text = []
    images: List[Dict[str, Any]] = []
    video_urls: List[str] = []

    for page_num, page in enumerate(doc):
        text = page.get_text()
        pages_text.append(text)
        video_urls.extend(_VIDEO_RE.findall(text))

        if not upload_images or len(images) >= MAX_IMAGES_PER_DOCUMENT:
            continue

        for img_info in page.get_images(full=True):
            if len(images) >= MAX_IMAGES_PER_DOCUMENT:
                break
            try:
                xref       = img_info[0]
                base_image = doc.extract_image(xref)
                img_bytes  = base_image["image"]
                ext        = base_image.get("ext", "png")
                if len(img_bytes) < 20_000:
                    continue
                try:
                    from PIL import Image as _PILImage
                    _pil = _PILImage.open(io.BytesIO(img_bytes))
                    if _pil.width < 100 or _pil.height < 100:
                        continue
                except Exception:
                    pass
                idx = len(images)
                url = _upload_image(img_bytes, ext, tenant_id, document_id, idx)
                images.append({"url": url, "page": page_num + 1, "index": idx, "ext": ext})
            except Exception as e:
                print(f"[media] PDF image error page {page_num}: {e}", flush=True)

    doc.close()
    return {
        "text":       "\n\n".join(p for p in pages_text if p.strip()),
        "images":     images,
        "video_urls": list(set(video_urls)),
    }


def _extract_docx(content: bytes, tenant_id: str, document_id: str, upload_images: bool = True) -> Dict[str, Any]:
    import zipfile
    from docx import Document

    doc  = Document(io.BytesIO(content))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    video_urls = list(set(_VIDEO_RE.findall(text)))

    images: List[Dict[str, Any]] = []
    if upload_images:
        with zipfile.ZipFile(io.BytesIO(content)) as z:
            for name in z.namelist():
                if not name.startswith("word/media/"):
                    continue
                ext = name.rsplit(".", 1)[-1].lower()
                if ext not in ("png", "jpg", "jpeg", "gif", "webp"):
                    continue
                try:
                    img_bytes = z.read(name)
                    if len(img_bytes) < 20_000:
                        continue
                    try:
                        from PIL import Image as _PILImage
                        _pil = _PILImage.open(io.BytesIO(img_bytes))
                        if _pil.width < 100 or _pil.height < 100:
                            continue
                    except Exception:
                        pass
                    if ext == "jpeg":
                        ext = "jpg"
                    idx = len(images)
                    url = _upload_image(img_bytes, ext, tenant_id, document_id, idx)
                    images.append({"url": url, "index": idx, "page": None, "ext": ext})
                except Exception as e:
                    print(f"[media] DOCX image error {name}: {e}", flush=True)

    return {"text": text, "images": images, "video_urls": video_urls}


# ── Faithful slide/page extraction ───────────────────────────────────────────

def _upload_slide_image(img_bytes: bytes, tenant_id: str, document_id: str, page: int) -> str:
    bucket = os.getenv("AWS_S3_BUCKET", "staffbot-docs")
    key    = f"{tenant_id}/{document_id}/slides/page_{page}.png"
    _s3().put_object(Bucket=bucket, Key=key, Body=img_bytes, ContentType="image/png")
    public_base = os.getenv("MINIO_PUBLIC_URL", "").rstrip("/")
    if public_base:
        return f"{public_base}/{key}"
    endpoint = os.getenv("AWS_ENDPOINT_URL", "http://localhost:9000")
    return f"{endpoint}/{bucket}/{key}"


def _first_nonempty_line(text: str) -> str:
    for line in text.splitlines():
        stripped = line.strip()
        if stripped:
            return stripped[:100]
    return ""


def _extract_pdf_faithful(content: bytes, tenant_id: str, document_id: str) -> List[Dict[str, Any]]:
    """Render each PDF page as PNG + extract text. Returns one dict per page."""
    import fitz

    doc    = fitz.open(stream=content, filetype="pdf")
    slides = []

    for page_num, page in enumerate(doc):
        text = page.get_text().strip()
        title = _first_nonempty_line(text) or f"Página {page_num + 1}"

        # Render at 150 DPI (scale factor ≈ 2.08 for 72 dpi base)
        mat      = fitz.Matrix(150 / 72, 150 / 72)
        pix      = page.get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes("png")

        image_url = None
        try:
            image_url = _upload_slide_image(img_bytes, tenant_id, document_id, page_num + 1)
        except Exception as e:
            print(f"[media] slide image upload error page {page_num + 1}: {e}", flush=True)

        slides.append({
            "page":         page_num + 1,
            "title":        title,
            "text":         text,
            "notes":        "",
            "image_url":    image_url,
            "source_format": "pdf",
        })

    doc.close()
    return slides


def _extract_pptx_faithful(content: bytes, tenant_id: str, document_id: str) -> List[Dict[str, Any]]:
    """Extract one entry per slide using python-pptx. Renders page images from embedded pics."""
    from pptx import Presentation
    from pptx.util import Pt
    import zipfile

    prs    = Presentation(io.BytesIO(content))
    slides = []

    # Build map: slide index → embedded media files
    # PPTX zip has ppt/slides/slide{n}.xml; media in ppt/media/
    # We correlate via relationship IDs per slide
    with zipfile.ZipFile(io.BytesIO(content)) as z:
        media_files = {name.split("/")[-1]: name for name in z.namelist() if name.startswith("ppt/media/")}

        for slide_idx, slide in enumerate(prs.slides):
            page_num = slide_idx + 1

            # ── Extract title ──────────────────────────────────────────
            title_text = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()

            # ── Extract all text ───────────────────────────────────────
            text_parts: List[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        text_parts.append(line)

            full_text = "\n".join(text_parts)
            if not title_text:
                title_text = _first_nonempty_line(full_text) or f"Slide {page_num}"

            # ── Extract speaker notes ──────────────────────────────────
            notes_text = ""
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                notes_text = tf.text.strip() if tf else ""

            # ── Extract first significant image from slide ─────────────
            image_url = None
            for rel in slide.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_bytes = rel.target_part.blob
                        if len(img_bytes) < 20_000:
                            continue
                        try:
                            from PIL import Image as _PILImage
                            _pil = _PILImage.open(io.BytesIO(img_bytes))
                            if _pil.width < 100 or _pil.height < 100:
                                continue
                        except Exception:
                            pass
                        # Determine extension from content type
                        ct  = rel.target_part.content_type
                        ext = "png" if "png" in ct else ("jpg" if "jpeg" in ct else "png")
                        key = f"{tenant_id}/{document_id}/slides/page_{page_num}_img.{ext}"
                        bucket = os.getenv("AWS_S3_BUCKET", "staffbot-docs")
                        _s3().put_object(Bucket=bucket, Key=key, Body=img_bytes, ContentType=f"image/{ext}")
                        public_base = os.getenv("MINIO_PUBLIC_URL", "").rstrip("/")
                        endpoint    = os.getenv("AWS_ENDPOINT_URL", "http://localhost:9000")
                        image_url   = f"{public_base}/{key}" if public_base else f"{endpoint}/{bucket}/{key}"
                        break  # one representative image per slide
                    except Exception as e:
                        print(f"[media] PPTX slide image error slide {page_num}: {e}", flush=True)

            slides.append({
                "page":          page_num,
                "title":         title_text,
                "text":          full_text,
                "notes":         notes_text,
                "image_url":     image_url,
                "source_format": "pptx",
            })

    return slides


def extract_slides_faithful(content: bytes, file_type: str, tenant_id: str, document_id: str) -> List[Dict[str, Any]]:
    """
    Extract one entry per page/slide, preserving order.
    Returns: [{ page, title, text, notes, image_url, source_format }]
    """
    if file_type == "pdf":
        return _extract_pdf_faithful(content, tenant_id, document_id)
    if file_type in ("pptx", "ppt", "odp"):
        return _extract_pptx_faithful(content, tenant_id, document_id)
    raise ValueError(f"Unsupported format for faithful extraction: {file_type}")


# ── Public entry point ────────────────────────────────────────────────────────

def extract_media(content: bytes, file_type: str, tenant_id: str, document_id: str, index_images: bool = True) -> Dict[str, Any]:
    """
    Extract text, images, and video URLs from a document.
    Returns { text, images: [{url, page?, index, ext}], video_urls: [str] }
    When index_images=False, images are not uploaded to MinIO and the images list is empty.
    """
    if file_type == "pdf":
        return _extract_pdf(content, tenant_id, document_id, upload_images=index_images)
    if file_type in ("docx", "doc"):
        return _extract_docx(content, tenant_id, document_id, upload_images=index_images)
    # txt / xlsx / other — text-only with video URL detection
    if file_type == "txt":
        text = content.decode("utf-8", errors="ignore")
    else:
        text = ""
    return {"text": text, "images": [], "video_urls": list(set(_VIDEO_RE.findall(text)))}
